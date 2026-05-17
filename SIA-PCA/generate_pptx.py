"""
Genera presentacion pptx con el seguimiento de los pasos del PCA
aplicado a europe.csv. Pocas slides, como pide la clase.
"""
import pandas as pd
import numpy as np
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

df = pd.read_csv("europe.csv")
X = df.drop(columns=["Country"])
X_std = (X - X.mean()) / X.std()
S = X.corr()
vals, vecs = np.linalg.eigh(S.values)
idx = np.argsort(vals)[::-1]
lambdas = vals[idx]
V = vecs[:, idx]
cols = list(S.columns)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

BLUE = RGBColor(0x4A, 0x6F, 0xC5)
DARK = RGBColor(0x33, 0x33, 0x33)

def add_title(slide, text, top=0.3):
    tx = slide.shapes.add_textbox(Inches(0.5), Inches(top), Inches(12.3), Inches(0.8))
    tf = tx.text_frame
    tf.text = text
    p = tf.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = DARK

def add_text(slide, text, left, top, width, height, size=14, bold=False, color=DARK, align=None):
    tx = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tx.text_frame
    tf.word_wrap = True
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = color
        if align is not None:
            p.alignment = align

# -------- Slide 1: Titulo --------
s = prs.slides.add_slide(prs.slide_layouts[6])
add_text(s, "Análisis de Componentes Principales", 0.5, 2.5, 12.3, 1, size=44, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
add_text(s, "Aplicación sobre europe.csv", 0.5, 3.6, 12.3, 0.6, size=22, align=PP_ALIGN.CENTER)
add_text(s, "Sistemas de Inteligencia Artificial - 2026", 0.5, 6.7, 12.3, 0.4, size=14, color=BLUE, align=PP_ALIGN.CENTER)

# -------- Slide 2: Pasos del procedimiento (slide 32) --------
s = prs.slides.add_slide(prs.slide_layouts[6])
add_title(s, "Procedimiento")
pasos = (
    "1.  Construir la matriz X (variables en columnas)\n"
    "2.  Estandarizar las variables\n"
    "3.  Calcular la matriz de correlaciones Sx\n"
    "4.  Calcular autovalores y autovectores (AyA)\n"
    "5.  Ordenar los autovalores de mayor a menor\n"
    "6.  Construir la matriz V con los autovectores\n"
    "7.  Calcular las nuevas variables Y como combinación lineal de las originales"
)
add_text(s, pasos, 1.2, 1.6, 11, 5, size=20)

# -------- Slide 3: Paso 1-2 - X y estandarizacion (boxplots) --------
s = prs.slides.add_slide(prs.slide_layouts[6])
add_title(s, "Pasos 1 y 2 — Matriz X y estandarización")
add_text(s, "28 países × 7 variables. Las escalas son muy distintas\n(Area ~10⁵ vs Inflation ~0-10) → se estandariza.",
         0.5, 1.2, 12.3, 1, size=14)
s.shapes.add_picture("boxplot_original.png", Inches(0.3), Inches(2.3), height=Inches(4.6))
s.shapes.add_picture("boxplot_estandarizado.png", Inches(6.8), Inches(2.3), height=Inches(4.6))

# -------- Slide 4: Paso 3 - matriz de correlaciones + ecuacion al costado --------
s = prs.slides.add_slide(prs.slide_layouts[6])
add_title(s, "Paso 3 — Matriz de correlaciones Sx")

# Matriz a la izquierda
rows = len(cols) + 1
table_shape = s.shapes.add_table(rows, len(cols)+1, Inches(0.3), Inches(1.4), Inches(8.5), Inches(4.2))
table = table_shape.table
table.cell(0, 0).text = ""
for j, c in enumerate(cols):
    cell = table.cell(0, j+1)
    cell.text = c[:9]
    for p in cell.text_frame.paragraphs:
        p.font.size = Pt(10); p.font.bold = True
for i, c in enumerate(cols):
    cell = table.cell(i+1, 0)
    cell.text = c[:9]
    for p in cell.text_frame.paragraphs:
        p.font.size = Pt(10); p.font.bold = True
    for j in range(len(cols)):
        cell = table.cell(i+1, j+1)
        val = S.iloc[i, j]
        cell.text = f"{val:+.2f}"
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(10)
            if i == j:
                p.font.bold = True

# Ecuaciones al costado derecho
add_text(s, "Fórmulas (slides 8 y 11)", 9.2, 1.4, 4, 0.4, size=14, bold=True, color=BLUE)
add_text(s,
    "Covarianza muestral:",
    9.2, 1.9, 4, 0.4, size=12, bold=True)
add_text(s,
    "s_ik = (1/(n-1)) · Σ (x_ji − x̄_i)(x_jk − x̄_k)",
    9.2, 2.3, 4, 0.6, size=12)
add_text(s,
    "Correlación muestral:",
    9.2, 3.0, 4, 0.4, size=12, bold=True)
add_text(s,
    "r_ik = s_ik / (√s_ii · √s_kk)",
    9.2, 3.4, 4, 0.5, size=12)
add_text(s,
    "Interpretación (slide 10):\n"
    "•  r > 0  → asociación positiva\n"
    "•  r < 0  → asociación negativa\n"
    "•  r ≈ 0  → independientes",
    9.2, 4.2, 4, 2, size=11)
add_text(s, "Valores |r| cercanos a 1 → redundancia (slide 16: si no hay correlación no tiene sentido PCA).",
         0.3, 5.8, 12, 0.5, size=11, color=BLUE)

# -------- Slide 5: Pasos 4-6 - autovalores y proporciones --------
s = prs.slides.add_slide(prs.slide_layouts[6])
add_title(s, "Pasos 4-6 — Autovalores, autovectores y proporción de varianza")
add_text(s, "λᵢ = varianza de Yᵢ  (slide 23)     |     Proporción = λᵢ / Σλ  (slide 31)",
         0.5, 1.2, 12.3, 0.4, size=13, bold=True, color=BLUE)

# tabla de autovalores
rows = len(lambdas) + 1
t = s.shapes.add_table(rows, 4, Inches(0.5), Inches(1.8), Inches(5.5), Inches(4.5)).table
headers = ["i", "λᵢ", "Prop.", "Acumulada"]
for j, h in enumerate(headers):
    t.cell(0, j).text = h
    for p in t.cell(0, j).text_frame.paragraphs:
        p.font.size = Pt(12); p.font.bold = True
acum = 0
for i, lv in enumerate(lambdas, 1):
    prop = lv / lambdas.sum()
    acum += prop
    vals_row = [str(i), f"{lv:.3f}", f"{prop*100:.2f}%", f"{acum*100:.2f}%"]
    for j, v in enumerate(vals_row):
        t.cell(i, j).text = v
        for p in t.cell(i, j).text_frame.paragraphs:
            p.font.size = Pt(11)
            if i == 1:
                p.font.bold = True; p.font.color.rgb = BLUE

# Primer autovector
add_text(s, "Primer autovector  v₁  (cargas de PC1):", 6.5, 1.8, 6, 0.4, size=13, bold=True)
loadings_text = "\n".join(f"{cols[i]:<13s}  {V[i,0]:+.4f}" for i in range(len(cols)))
add_text(s, loadings_text, 6.5, 2.3, 6, 4, size=13)
add_text(s, "Y₁ = Σ vᵢ₁ · Xᵢ_std    (slides 20, 22)", 6.5, 6.0, 6, 0.5, size=12, color=BLUE)

# -------- Slide 6: Paso 7 - Y1 grafico de barras + biplot --------
s = prs.slides.add_slide(prs.slide_layouts[6])
add_title(s, "Paso 7 — Componentes Y: índice (PC1) y biplot (PC1 vs PC2)")
s.shapes.add_picture("pc1_por_pais.png", Inches(0.2), Inches(1.4), height=Inches(5.6))
s.shapes.add_picture("biplot.png", Inches(6.8), Inches(1.4), height=Inches(5.6))

# -------- Slide 7: Conclusion --------
s = prs.slides.add_slide(prs.slide_layouts[6])
add_title(s, "Conclusión — Interpretación de PC1")
texto = (
    "•  PC1 captura el 46.1% de la variabilidad total; con PC1+PC2 llegamos al 63%.\n\n"
    "•  Cargas dominantes en PC1:  GDP (+0.50), Life.expect (+0.48), Pop.growth (+0.48)\n"
    "    en sentido positivo;  Inflation (−0.41) y Unemployment (−0.27) en sentido negativo.\n\n"
    "•  Países con Y₁ alto (Luxembourg, Switzerland, Norway) ↔ alto GDP, alta expectativa\n"
    "    de vida, baja inflación, bajo desempleo.\n"
    "    Países con Y₁ bajo (Ukraine, Bulgaria, Estonia) ↔ situación opuesta.\n\n"
    "•  PC1 actúa como un ÍNDICE (slide 38) que resume las 7 variables en una sola\n"
    "    dimensión interpretable como  «nivel de desarrollo socioeconómico».\n\n"
    "•  La redundancia entre variables correlacionadas (GDP↔Life.expect, GDP↔Pop.growth,\n"
    "    Inflation↔Life.expect) fue eliminada: las componentes Y son no correlacionadas\n"
    "    entre sí (slide 14)."
)
add_text(s, texto, 0.7, 1.5, 12, 6, size=15)

prs.save("PCA_europe_seguimiento.pptx")
print("Guardado: PCA_europe_seguimiento.pptx")
